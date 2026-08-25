import os
import io
import json
from datetime import datetime
import pytz
import pandas as pd
import streamlit as st
import streamlit.components.v1 as components
from PIL import Image
from groq import Groq
import pypdf
from docx import Document
from pptx import Presentation

# --- CONFIGURACION DE PAGINA ---
st.set_page_config(
    page_title="Plataforma de Analitica & Materiales IA",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- RUTAS DE ALMACENAMIENTO PERSISTENTE ---
DIR_BASE = "almacen_datos"
DIR_DATASETS = os.path.join(DIR_BASE, "datasets")
DIR_MATERIAL = os.path.join(DIR_BASE, "material")
DIR_INFORMES = os.path.join(DIR_BASE, "informes_guardados")
DIR_USUARIOS = os.path.join(DIR_BASE, "usuarios_guardados")
DIR_AVATARS = os.path.join(DIR_BASE, "avatares")
FILE_DB = os.path.join(DIR_BASE, "base_datos.json")

for d in [DIR_BASE, DIR_DATASETS, DIR_MATERIAL, DIR_INFORMES, DIR_USUARIOS, DIR_AVATARS]:
    os.makedirs(d, exist_ok=True)[cite: 7, 8, 9]

MODELO_WHISPER = "whisper-large-v3"[cite: 7, 8, 9]

# --- ESTILOS VISUALES: LAVANDA CLARO & MORADO ---
st.markdown("""
<style>
    html, body, [class*="css"], .stApp {
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
    }

    /* Fondo principal: Lavanda claro */
    .stApp {
        background: #f3e8ff !important;
        color: #1e1b4b !important;
    }
    
    h1, h2, h3, h4, h5, h6 {
        color: #3b0764 !important;
        font-weight: 800;
        letter-spacing: -0.5px;
        font-family: 'Segoe UI', system-ui, -apple-system, sans-serif !important;
    }
    
    /* Barra lateral */
    section[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #6214c7 0%, #4c1d95 100%) !important;
        border-right: 2px solid #3b0764 !important;
    }
    section[data-testid="stSidebar"] * {
        color: #f5f3ff !important;
        font-weight: 600;
    }
    section[data-testid="stSidebar"] h1, 
    section[data-testid="stSidebar"] h2, 
    section[data-testid="stSidebar"] h3 {
        color: #ffffff !important;
    }

    /* Tarjetas y Contenedores */
    .modern-card {
        background: #ffffff !important;
        border: 2px solid #d8b4fe !important;
        border-radius: 14px;
        padding: 22px;
        margin-bottom: 20px;
        box-shadow: 0 4px 14px rgba(91, 33, 182, 0.08);
    }

    /* Desplegables (Expanders) */
    div[data-testid="stExpander"] {
        background-color: #ffffff !important;
        border: 2px solid #d8b4fe !important;
        border-radius: 10px !important;
        margin-bottom: 14px !important;
    }
    div[data-testid="stExpander"] summary {
        background: linear-gradient(135deg, #7c3aed 0%, #6d28d9 100%) !important;
        color: #ffffff !important;
        border-radius: 8px !important;
        font-weight: 750 !important;
    }
    div[data-testid="stExpander"] summary * {
        color: #ffffff !important;
        font-weight: 750 !important;
    }
    div[data-testid="stExpander"] div[role="region"] {
        background-color: #ffffff !important;
        color: #1e1b4b !important;
    }

    /* Inputs y Textareas */
    input[type="text"], 
    input[type="password"], 
    textarea {
        background-color: #faf5ff !important;
        color: #1e1b4b !important;
        border: 2px solid #c084fc !important;
        border-radius: 8px !important;
        font-size: 0.95rem !important;
        font-weight: 600 !important;
    }
    input[type="text"]:focus, 
    input[type="password"]:focus, 
    textarea:focus {
        background-color: #ffffff !important;
        color: #000000 !important;
        border-color: #7c3aed !important;
        box-shadow: 0 0 0 3px rgba(124, 58, 237, 0.25) !important;
    }

    div[data-baseweb="input"] button,
    div[data-baseweb="input"] div {
        background-color: #faf5ff !important;
        border: none !important;
        color: #1e1b4b !important;
    }
    div[data-baseweb="input"] svg {
        fill: #1e1b4b !important;
        color: #1e1b4b !important;
    }

    div[data-baseweb="select"] > div {
        background-color: #faf5ff !important;
        color: #1e1b4b !important;
        border: 2px solid #c084fc !important;
        border-radius: 8px !important;
    }
    div[data-baseweb="select"] * {
        background-color: #faf5ff !important;
        color: #1e1b4b !important;
        font-weight: 600 !important;
    }
    div[data-baseweb="select"] svg {
        fill: #1e1b4b !important;
        color: #1e1b4b !important;
    }

    /* Checkboxes */
    div[data-testid="stCheckbox"] span[role="checkbox"] {
        background-color: #faf5ff !important;
        border: 2px solid #a855f7 !important;
        border-radius: 5px !important;
    }
    div[data-testid="stCheckbox"] span[role="checkbox"][aria-checked="true"] {
        background-color: #7c3aed !important;
        border-color: #6d28d9 !important;
    }
    div[data-testid="stCheckbox"] svg {
        color: #ffffff !important;
        fill: #ffffff !important;
    }
    div[data-testid="stCheckbox"] label span {
        color: #1e1b4b !important;
        font-weight: 600 !important;
    }

    /* Subida de archivos */
    div[data-testid="stFileUploader"] {
        background-color: #faf5ff !important;
        border: 2px dashed #a855f7 !important;
        border-radius: 12px !important;
        padding: 14px !important;
    }
    div[data-testid="stFileUploader"] section {
        background-color: #faf5ff !important;
    }
    div[data-testid="stFileUploader"] section * {
        color: #1e1b4b !important;
    }
    div[data-testid="stFileUploader"] button {
        background: linear-gradient(135deg, #7c3aed 0%, #6d28d9 100%) !important;
        border: 2px solid #6d28d9 !important;
        border-radius: 8px !important;
        font-weight: 750 !important;
    }
    div[data-testid="stFileUploader"] button p {
        color: #ffffff !important;
    }

    /* Botones de Descarga */
    div[data-testid="stDownloadButton"]>button {
        background: #ede9fe !important;
        color: #4c1d95 !important;
        border: 2px solid #a855f7 !important;
        border-radius: 9px !important;
        padding: 9px 20px !important;
        font-weight: 800 !important;
        font-size: 0.95rem !important;
        box-shadow: 0 3px 8px rgba(76, 29, 149, 0.12) !important;
        transition: all 0.2s ease !important;
    }
    div[data-testid="stDownloadButton"]>button:hover {
        background: #ddd6fe !important;
        color: #1e1b4b !important;
        border-color: #7c3aed !important;
        transform: translateY(-1px);
    }
    div[data-testid="stDownloadButton"]>button p {
        color: #4c1d95 !important;
        font-weight: 800 !important;
    }

    /* Botones Principales */
    .stButton>button, 
    div[data-testid="stFormSubmitButton"]>button {
        background: linear-gradient(135deg, #7c3aed 0%, #6214c7 100%) !important;
        color: #ffffff !important;
        border: none !important;
        border-radius: 9px !important;
        padding: 9px 22px !important;
        font-weight: 800 !important;
        font-size: 0.95rem !important;
        box-shadow: 0 4px 12px rgba(109, 40, 217, 0.3) !important;
        transition: all 0.2s ease !important;
    }
    .stButton>button:hover, 
    div[data-testid="stFormSubmitButton"]>button:hover {
        background: linear-gradient(135deg, #8b5cf6 0%, #7c3aed 100%) !important;
        box-shadow: 0 6px 16px rgba(109, 40, 217, 0.45) !important;
        transform: translateY(-1px);
    }
    .stButton>button p,
    div[data-testid="stFormSubmitButton"]>button p {
        color: #ffffff !important;
        font-weight: 800 !important;
    }

    /* Pestañas (Tabs) */
    [data-testid="stTabs"] [data-baseweb="tab-list"] {
        gap: 8px !important;
        border-bottom: 2.5px solid #c084fc !important;
        background: transparent !important;
    }
    [data-testid="stTabs"] button[role="tab"] {
        background-color: #ede9fe !important;
        border: 1.5px solid #d8b4fe !important;
        border-bottom: none !important;
        border-radius: 10px 10px 0 0 !important;
        padding: 10px 20px !important;
    }
    [data-testid="stTabs"] button[role="tab"] * {
        color: #5b21b6 !important;
        -webkit-text-fill-color: #5b21b6 !important;
        font-weight: 800 !important;
        font-size: 1.02rem !important;
    }
    [data-testid="stTabs"] button[role="tab"][aria-selected="true"] {
        background-color: #6214c7 !important;
        border: 1.5px solid #6214c7 !important;
    }
    [data-testid="stTabs"] button[role="tab"][aria-selected="true"] * {
        color: #ffffff !important;
        -webkit-text-fill-color: #ffffff !important;
        font-weight: 900 !important;
    }

    /* Tablas */
    table {
        width: 100% !important;
        border-collapse: collapse !important;
        margin: 16px 0 !important;
        background-color: #ffffff !important;
        border-radius: 8px !important;
        overflow: hidden !important;
        border: 2px solid #d8b4fe !important;
    }
    th {
        background-color: #ede9fe !important;
        color: #3b0764 !important;
        font-weight: 800 !important;
        padding: 12px 14px !important;
        border: 1px solid #d8b4fe !important;
        text-align: left !important;
    }
    td {
        padding: 10px 14px !important;
        border: 1px solid #e9d5ff !important;
        color: #1e1b4b !important;
        font-size: 0.95rem !important;
    }
    tr:nth-child(even) {
        background-color: #faf5ff !important;
    }

    .highlight-tag {
        background: #ede9fe;
        border: 1.5px solid #c084fc;
        padding: 3px 10px;
        border-radius: 6px;
        font-weight: 700;
        color: #4c1d95;
        font-size: 0.92rem;
        display: inline-block;
    }

    .badge-role {
        display: inline-block;
        background: #7c3aed;
        color: #ffffff !important;
        padding: 4px 12px;
        border-radius: 14px;
        font-size: 0.82rem;
        font-weight: 800;
    }

    .login-container-lavender {
        background: #ffffff !important;
        border: 2px solid #c084fc !important;
        border-radius: 18px;
        padding: 38px 34px;
        box-shadow: 0 12px 32px rgba(109, 40, 217, 0.15);
        margin-top: 40px;
    }
</style>
""", unsafe_allow_html=True)

# --- RELOJ DUAL ---
def renderizar_reloj_chile():
    html_reloj = """
    <div style="background: linear-gradient(135deg, #7c3aed 0%, #6214c7 100%); border: 2px solid #5b21b6; border-radius: 16px; padding: 12px 20px; display: flex; align-items: center; justify-content: space-between; gap: 16px; box-shadow: 0 6px 18px rgba(109, 40, 217, 0.35); max-width: 480px; margin-left: auto; margin-bottom: 15px; font-family: 'Segoe UI', system-ui, sans-serif;">
        <div style="position: relative; width: 60px; height: 60px;">
            <svg id="analog-clock" width="60" height="60" viewBox="0 0 100 100">
                <circle cx="50" cy="50" r="46" fill="#faf5ff" stroke="#4c1d95" stroke-width="4"/>
                <line x1="50" y1="10" x2="50" y2="16" stroke="#4c1d95" stroke-width="3" stroke-linecap="round"/>
                <line x1="90" y1="50" x2="84" y2="50" stroke="#4c1d95" stroke-width="3" stroke-linecap="round"/>
                <line x1="50" y1="90" x2="50" y2="84" stroke="#4c1d95" stroke-width="3" stroke-linecap="round"/>
                <line x1="10" y1="50" x2="16" y2="50" stroke="#4c1d95" stroke-width="3" stroke-linecap="round"/>
                <line id="hour-hand" x1="50" y1="50" x2="50" y2="28" stroke="#1e1b4b" stroke-width="5" stroke-linecap="round"/>
                <line id="min-hand" x1="50" y1="50" x2="50" y2="18" stroke="#5b21b6" stroke-width="3.5" stroke-linecap="round"/>
                <line id="sec-hand" x1="50" y1="50" x2="50" y2="14" stroke="#c084fc" stroke-width="2" stroke-linecap="round"/>
                <circle cx="50" cy="50" r="4" fill="#1e1b4b"/>
            </svg>
        </div>
        <div style="display: flex; flex-direction: column; justify-content: center;">
            <div style="font-size: 0.76rem; font-weight: 800; color: #ede9fe; text-transform: uppercase; letter-spacing: 0.8px;">🇨🇱 Hora Oficial de Chile</div>
            <div id="digital-clock" style="font-size: 1.65rem; font-weight: 900; color: #ffffff; line-height: 1.1; text-shadow: 0 2px 4px rgba(0,0,0,0.2);">--:--:--</div>
            <div id="digital-date" style="font-size: 0.82rem; font-weight: 700; color: #ede9fe; margin-top: 2px;">Cargando fecha...</div>
        </div>
    </div>
    <script>
        function actualizarReloj() {
            const opciones = { timeZone: 'America/Santiago', hour12: false, hour: '2-digit', minute: '2-digit', second: '2-digit' };
            const opcionesFecha = { timeZone: 'America/Santiago', weekday: 'long', year: 'numeric', month: 'long', day: 'numeric' };
            const ahora = new Date();
            const formateadorHora = new Intl.DateTimeFormat('es-CL', opciones);
            const formateadorFecha = new Intl.DateTimeFormat('es-CL', opcionesFecha);
            const partesHora = formateadorHora.formatToParts(ahora);
            let h = 0, m = 0, s = 0;
            partesHora.forEach(p => {
                if (p.type === 'hour') h = parseInt(p.value);
                if (p.type === 'minute') m = parseInt(p.value);
                if (p.type === 'second') s = parseInt(p.value);
            });
            const hStr = h.toString().padStart(2, '0');
            const mStr = m.toString().padStart(2, '0');
            const sStr = s.toString().padStart(2, '0');
            document.getElementById('digital-clock').textContent = `${hStr}:${mStr}:${sStr}`;
            let fechaStr = formateadorFecha.format(ahora);
            fechaStr = fechaStr.charAt(0).toUpperCase() + fechaStr.slice(1);
            document.getElementById('digital-date').textContent = fechaStr;
            const secDeg = s * 6;
            const minDeg = m * 6 + s * 0.1;
            const hourDeg = (h % 12) * 30 + m * 0.5;
            function rotar(elemId, deg) {
                const el = document.getElementById(elemId);
                if (el) el.setAttribute('transform', `rotate(${deg} 50 50)`);
            }
            rotar('sec-hand', secDeg);
            rotar('min-hand', minDeg);
            rotar('hour-hand', hourDeg);
        }
        setInterval(actualizarReloj, 1000);
        actualizarReloj();
    </script>
    """
    components.html(html_reloj, height=105)

# --- GESTOR DE PERSISTENCIA Y RESTAURACION TOTAL ---
def sincronizar_usuarios_fisicos(usuarios_dict):
    for u_k, u_data in usuarios_dict.items():
        ruta_u_json = os.path.join(DIR_USUARIOS, f"{u_k}.json")
        with open(ruta_u_json, "w", encoding="utf-8") as fu:
            json.dump(u_data, fu, ensure_ascii=False, indent=2)

def restaurar_repositorio_local(data):
    if os.path.exists(DIR_USUARIOS):
        for arch in os.listdir(DIR_USUARIOS):
            if arch.endswith(".json"):
                u_login = arch.rsplit(".json", 1)[0]
                ruta_uj = os.path.join(DIR_USUARIOS, arch)
                try:
                    with open(ruta_uj, "r", encoding="utf-8") as fuj:
                        u_info = json.load(fuj)
                        if u_login not in data["usuarios"]:
                            data["usuarios"][u_login] = u_info
                except Exception:
                    pass

    if os.path.exists(DIR_DATASETS):
        for arch in os.listdir(DIR_DATASETS):
            ruta_f = os.path.join(DIR_DATASETS, arch)
            if os.path.isfile(ruta_f):
                partes = arch.split("_", 1)
                titulo_ds = partes[1].rsplit(".", 1)[0] if len(partes) > 1 else arch
                if titulo_ds not in data["datasets"]:
                    data["datasets"][titulo_ds] = {
                        "titulo": titulo_ds,
                        "autor": "Archivo Restaurado",
                        "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "ruta": ruta_f
                    }

    rutas_mat = {item.get("ruta") for item in data.get("material", []) if item.get("ruta")}
    if os.path.exists(DIR_MATERIAL):
        for arch in os.listdir(DIR_MATERIAL):
            ruta_f = os.path.join(DIR_MATERIAL, arch)
            if ruta_f not in rutas_mat and os.path.isfile(ruta_f):
                partes = arch.split("_", 1)
                nombre_orig = partes[1] if len(partes) > 1 else arch
                ext = arch.split(".")[-1].lower()
                data["material"].append({
                    "titulo": f"Material Restaurado ({nombre_orig})",
                    "nombre_original": nombre_orig,
                    "tipo": ext,
                    "autor": "Sistema / Restauración",
                    "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                    "ruta": ruta_f,
                    "resumen": "Archivo físico recuperado automáticamente del almacenamiento persistente."
                })

    if os.path.exists(DIR_INFORMES):
        titulos_inf = {inf.get("titulo") for inf in data.get("informes", [])}
        for arch in os.listdir(DIR_INFORMES):
            if arch.endswith(".json"):
                ruta_inf = os.path.join(DIR_INFORMES, arch)
                try:
                    with open(ruta_inf, "r", encoding="utf-8") as fi:
                        inf_obj = json.load(fi)
                        if inf_obj.get("titulo") not in titulos_inf:
                            data["informes"].append(inf_obj)
                except Exception:
                    pass
    return data

def cargar_estado():
    data_defecto = {
        "usuarios": {
            "Francesca Fellay": {
                "nombre": "Francesca Fellay",
                "rol": "Admin",
                "pin": "1953",
                "permiso_editar": True,
                "permiso_eliminar": True,
                "avatar": ""
            },
            "gerente": {
                "nombre": "Gerencia General",
                "rol": "Admin",
                "pin": "gerente",
                "permiso_editar": True,
                "permiso_eliminar": True,
                "avatar": ""
            }
        },
        "datasets": {},
        "material": [],
        "informes": []
    }
    
    if not os.path.exists(FILE_DB):
        data = restaurar_repositorio_local(data_defecto)
        sincronizar_usuarios_fisicos(data["usuarios"])
        guardar_estado(data)
        return data
        
    with open(FILE_DB, "r", encoding="utf-8") as f:
        try:
            data = json.load(f)
        except Exception:
            data = data_defecto

    if "material" not in data:
        data["material"] = data.pop("intervencion", [])
    if "datasets" not in data:
        data["datasets"] = {}
    if "informes" not in data:
        data["informes"] = []
    if "usuarios" not in data:
        data["usuarios"] = {}

    data["usuarios"]["Francesca Fellay"] = {
        "nombre": "Francesca Fellay",
        "rol": "Admin",
        "pin": "1953",
        "permiso_editar": True,
        "permiso_eliminar": True,
        "avatar": data["usuarios"].get("Francesca Fellay", {}).get("avatar", "")
    }
    data["usuarios"]["gerente"] = {
        "nombre": "Gerencia General",
        "rol": "Admin",
        "pin": "gerente",
        "permiso_editar": True,
        "permiso_eliminar": True,
        "avatar": ""
    }
    
    data = restaurar_repositorio_local(data)
    sincronizar_usuarios_fisicos(data["usuarios"])
    guardar_estado(data)
    return data

def guardar_estado(data):
    sincronizar_usuarios_fisicos(data.get("usuarios", {}))
    with open(FILE_DB, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

db = cargar_estado()

# --- CLIENTE GROQ IA ---
def obtener_api_key_groq():
    if "GROQ_API_KEY" in st.secrets:
        return st.secrets["GROQ_API_KEY"]
    return os.environ.get("GROQ_API_KEY", "")

API_KEY_GROQ = obtener_api_key_groq()

def obtener_cliente_ia():
    if not API_KEY_GROQ:
        return None
    try:
        return Groq(api_key=API_KEY_GROQ)
    except Exception:
        return None

def ejecutar_chat_groq(client, prompt_sistema, prompt_usuario):
    modelos_candidatos = [
        "llama-3.3-70b-versatile",
        "llama-3.1-70b-versatile",
        "llama3-70b-8192",
        "mixtral-8x7b-32768",
        "gemma2-9b-it"
    ]
    try:
        lista_api = [m.id for m in client.models.list().data if "whisper" not in m.id.lower() and "guard" not in m.id.lower()]
    except Exception:
        lista_api = []

    candidatos = [m for m in modelos_candidatos if m in lista_api] + lista_api + modelos_candidatos
    candidatos = list(dict.fromkeys(candidatos))

    ultimo_error = None
    for model_id in candidatos:
        try:
            resp = client.chat.completions.create(
                model=model_id,
                messages=[
                    {"role": "system", "content": prompt_sistema},
                    {"role": "user", "content": prompt_usuario}
                ],
                temperature=0.3
            )
            return resp.choices[0].message.content
        except Exception as e:
            ultimo_error = e
            continue
    raise Exception(f"No fue posible conectar con los modelos de Groq. Detalle: {ultimo_error}")

# --- EXTRACCION DE TEXTO ---
def extraer_texto_archivo(ruta, extension):
    texto = ""
    try:
        if extension == "pdf":
            lector = pypdf.PdfReader(ruta)
            for pag in lector.pages:
                t = pag.extract_text()
                if t:
                    texto += t + "\n"
        elif extension == "docx":
            doc = Document(ruta)
            texto = "\n".join([p.text for p in doc.paragraphs])
        elif extension == "pptx":
            prs = Presentation(ruta)
            for diap in prs.slides:
                for forma in diap.shapes:
                    if hasattr(forma, "text"):
                        texto += forma.text + "\n"
        elif extension in ["xlsx", "xls"]:
            df_tmp = pd.read_excel(ruta)
            texto = f"Estadisticas:\n{df_tmp.describe(include='all').to_string()}\nPrimeras filas:\n{df_tmp.head(10).to_string()}"
    except Exception as e:
        texto = f"Error al extraer texto: {e}"
    return texto[:8000]

# --- CONTROL DE ACCESO (LOGIN) ---
if "autenticado" not in st.session_state:
    st.session_state.autenticado = False
    st.session_state.usuario_clave = None

if not st.session_state.autenticado:
    c1, c2, c3 = st.columns([1, 1.4, 1])
    with c2:
        st.markdown("""
        <div class='login-container-lavender'>
            <div style='text-align: center; margin-bottom: 22px;'>
                <div style='font-size: 3rem; margin-bottom: 4px;'>⚡</div>
                <h2 style='margin: 0; color: #3b0764 !important; font-size: 1.8rem;'>Acceso a la Plataforma</h2>
                <p style='margin: 4px 0 0 0; color: #6b21a8; font-size: 0.95rem; font-weight: 700;'>Gestión & Analítica con Groq IA</p>
            </div>
        """, unsafe_allow_html=True)
        
        with st.form("login_form"):
            st.markdown("<p style='color:#3b0764; font-weight:750; margin-bottom:4px;'>Usuario:</p>", unsafe_allow_html=True)
            u_in = st.text_input("Usuario", label_visibility="collapsed")
            
            st.markdown("<p style='color:#3b0764; font-weight:750; margin-bottom:4px; margin-top:12px;'>PIN / Contraseña:</p>", unsafe_allow_html=True)
            p_in = st.text_input("PIN", type="password", label_visibility="collapsed")
            
            st.markdown("<div style='margin-top: 18px;'></div>", unsafe_allow_html=True)
            if st.form_submit_button("Iniciar Sesión", use_container_width=True):
                u_clean = u_in.strip()
                if u_clean in db["usuarios"] and db["usuarios"][u_clean]["pin"] == p_in:
                    st.session_state.autenticado = True
                    st.session_state.usuario_clave = u_clean
                    st.rerun()
                else:
                    st.error("Credenciales incorrectas. Verifique el usuario y el PIN ingresado.")
                    
        st.markdown("</div>", unsafe_allow_html=True)
    st.stop()

# --- DATOS DEL USUARIO ACTUAL ---
usr_key = st.session_state.usuario_clave
usr = db["usuarios"].get(usr_key, {"nombre": usr_key, "rol": "Admin", "permiso_editar": True, "permiso_eliminar": True, "avatar": ""})
es_admin = (usr.get("rol") == "Admin") or (usr_key in ["Francesca Fellay", "gerente"])

# --- BARRA LATERAL ---
st.sidebar.markdown("### 👤 Mi Perfil")

avatar_path = usr.get("avatar", "")
if avatar_path and os.path.exists(avatar_path):
    st.sidebar.image(avatar_path, width=105)
else:
    st.sidebar.markdown("""<div style='width:80px; height:80px; border-radius:50%; background:#ede9fe; border:2px solid #a855f7; display:flex; align-items:center; justify-content:center; font-size:2.2rem; color:#4c1d95; margin-bottom:10px;'>👤</div>""", unsafe_allow_html=True)

st.sidebar.markdown(f"**{usr.get('nombre')}**")
st.sidebar.markdown(f"<span class='badge-role'>{usr.get('rol')}</span>", unsafe_allow_html=True)

with st.sidebar.expander("📷 Cambiar foto de perfil"):
    nueva_foto = st.file_uploader("Subir imagen (JPG, PNG):", type=["jpg", "jpeg", "png"], key="upload_avatar")
    if nueva_foto is not None:
        if st.button("Guardar Foto"):
            ruta_avatar = os.path.join(DIR_AVATARS, f"{usr_key}_avatar.png")
            with open(ruta_avatar, "wb") as f_av:
                f_av.write(nueva_foto.getbuffer())
            db["usuarios"][usr_key]["avatar"] = ruta_avatar
            guardar_estado(db)
            st.success("Foto actualizada.")
            st.rerun()

st.sidebar.markdown("---")
if st.sidebar.button("🚪 Cerrar Sesión", use_container_width=True):
    st.session_state.autenticado = False
    st.session_state.usuario_clave = None
    st.rerun()

# --- ENCABEZADO SUPERIOR CON RELOJ ---
col_head_title, col_head_clock = st.columns([1.2, 1])
with col_head_title:
    st.markdown("<h1 style='color:#3b0764 !important; margin:0;'>⚡ Centro de Gestión & Analítica</h1>", unsafe_allow_html=True)
    st.markdown("<p style='color:#6b21a8; margin-top:4px; font-weight:600;'>Plataforma colaborativa multi-formato con procesamiento inteligente mediante Groq IA.</p>", unsafe_allow_html=True)
with col_head_clock:
    renderizar_reloj_chile()

# --- NAVEGACION POR PESTAÑAS ---
titulos_pestanas = ["📊 Datasets & Archivos", "📂 Material", "📄 Informes Compartidos"]
if es_admin:
    titulos_pestanas.append("👥 Gestión de Usuarios")

pestanas_principales = st.tabs(titulos_pestanas)

# ==========================================
# 1. DATASETS EXCEL
# ==========================================
with pestanas_principales[0]:
    st.markdown("""<div class='modern-card'><h3 style='margin:0;'>📊 Repositorio de Datasets</h3><p style='margin:0; color:#4c1d95; font-weight:600;'>Suba, visualice y edite hojas de cálculo en ventanas independientes.</p></div>""", unsafe_allow_html=True)
    
    with st.expander("➕ Cargar Nuevo Dataset Excel", expanded=True):
        c_tit, c_arc = st.columns([1, 1])
        with c_tit:
            t_data = st.text_input("Título del dataset:")
        with c_arc:
            f_data = st.file_uploader("Archivo Excel (.xlsx, .xls):", type=["xlsx", "xls"])
            
        if st.button("Guardar Dataset"):
            if not t_data.strip() or f_data is None:
                st.error("Complete el título y seleccione un archivo.")
            else:
                nom_arc = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{f_data.name}"
                ruta_dest = os.path.join(DIR_DATASETS, nom_arc)
                with open(ruta_dest, "wb") as f:
                    f.write(f_data.getbuffer())
                
                db["datasets"][t_data] = {
                    "titulo": t_data,
                    "autor": usr["nombre"],
                    "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                    "ruta": ruta_dest
                }
                guardar_estado(db)
                st.success("Dataset guardado con éxito.")
                st.rerun()

    st.markdown("---")
    if not db["datasets"]:
        st.info("No hay datasets subidos actualmente.")
    else:
        titulos_ds = list(db["datasets"].keys())
        pestanas_ds = st.tabs(titulos_ds)
        
        for idx_ds, tab_ds in enumerate(pestanas_ds):
            t_act = titulos_ds[idx_ds]
            info_ds = db["datasets"][t_act]
            
            with tab_ds:
                st.markdown(f"### 📋 {info_ds['titulo']}")
                st.caption(f"Subido por: **{info_ds['autor']}** | Fecha: {info_ds['fecha']}")
                
                if os.path.exists(info_ds["ruta"]):
                    df_actual = pd.read_excel(info_ds["ruta"])
                    
                    if usr.get("permiso_editar") or es_admin:
                        st.markdown("**✏️ Editor de Datos en Vivo:**")
                        df_edit = st.data_editor(df_actual, key=f"d_edit_{t_act}", use_container_width=True)
                        if st.button("💾 Guardar Cambios en Excel", key=f"s_df_{t_act}"):
                            df_edit.to_excel(info_ds["ruta"], index=False)
                            st.success("Datos actualizados.")
                            st.rerun()
                    else:
                        st.dataframe(df_actual, use_container_width=True)
                    
                    if usr.get("permiso_eliminar") or es_admin:
                        st.markdown("---")
                        if st.button("🗑️ Eliminar Dataset", key=f"del_ds_{t_act}", type="secondary"):
                            if os.path.exists(info_ds["ruta"]):
                                os.remove(info_ds["ruta"])
                            del db["datasets"][t_act]
                            guardar_estado(db)
                            st.success("Dataset eliminado.")
                            st.rerun()
                else:
                    st.error("Archivo físico no encontrado.")

# ==========================================
# 2. MODULO MATERIAL MULTI-FORMATO
# ==========================================
with pestanas_principales[1]:
    st.markdown("""<div class='modern-card'><h3 style='margin:0;'>📂 Módulo de Material</h3><p style='margin:0; color:#4c1d95; font-weight:600;'>Soporte y resúmenes automáticos con Groq para Documentos, Imágenes y Audios (M4A/MP3/WAV).</p></div>""", unsafe_allow_html=True)
    
    with st.form("form_material"):
        tit_mat = st.text_input("Título descriptivo del material:")
        archivos = st.file_uploader(
            "Cargar archivos:",
            type=["xlsx", "xls", "docx", "pdf", "pptx", "jpg", "jpeg", "png", "mp4", "m4a", "mp3", "wav"],
            accept_multiple_files=True
        )
        if st.form_submit_button("Subir y Procesar"):
            if not tit_mat.strip() or not archivos:
                st.error("Complete el título y cargue al menos un archivo.")
            else:
                client = obtener_cliente_ia()
                for a in archivos:
                    ext = a.name.split(".")[-1].lower()
                    nom_dest = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{a.name}"
                    ruta_guardada = os.path.join(DIR_MATERIAL, nom_dest)
                    
                    with open(ruta_guardada, "wb") as f:
                        f.write(a.getbuffer())
                    
                    resumen_txt = "Archivo guardado."
                    if client:
                        with st.spinner(f"Analizando {a.name} con Groq IA..."):
                            try:
                                if ext in ["m4a", "mp3", "wav"]:
                                    with open(ruta_guardada, "rb") as f_aud:
                                        audio_bytes = f_aud.read()
                                    audio_buffer = io.BytesIO(audio_bytes)
                                    audio_buffer.name = f"temp_audio.{ext}"
                                    
                                    transcripcion = client.audio.transcriptions.create(
                                        model=MODELO_WHISPER,
                                        file=audio_buffer,
                                        language="es"
                                    )
                                    texto_audio = transcripcion.text
                                    
                                    p_sys = "Eres un especialista senior en diagnóstico y análisis de materiales. Sintetiza con precisión en español latinoamericano."
                                    p_user = f"A partir de la siguiente transcripción de audio ({a.name}), sintetiza los puntos clave tratados, acuerdos y diagnóstico. Usa párrafos claros o viñetas Markdown estándar, nunca diagramas ASCII:\n\n{texto_audio}"
                                    resumen_txt = ejecutar_chat_groq(client, p_sys, p_user)
                                elif ext in ["pdf", "docx", "pptx", "xlsx", "xls"]:
                                    t_doc = extraer_texto_archivo(ruta_guardada, ext)
                                    p_sys = "Eres un especialista senior en análisis documental y materiales."
                                    p_user = f"Elabora un resumen y diagnóstico clave del siguiente documento ({a.name}):\n\n{t_doc}\n\nUsa viñetas Markdown estándar con conceptos clave en negrita."
                                    resumen_txt = ejecutar_chat_groq(client, p_sys, p_user)
                                elif ext in ["jpg", "jpeg", "png"]:
                                    resumen_txt = f"Imagen registrada ({a.name}). Visualización disponible en la pestaña multimedia."
                                elif ext == "mp4":
                                    resumen_txt = f"Video registrado ({a.name}). Reproducción multimedia disponible."
                            except Exception as e:
                                resumen_txt = f"Archivo guardado. Diagnóstico no generado: {e}"
                    
                    db["material"].append({
                        "titulo": tit_mat,
                        "nombre_original": a.name,
                        "tipo": ext,
                        "autor": usr["nombre"],
                        "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "ruta": ruta_guardada,
                        "resumen": resumen_txt
                    })
                guardar_estado(db)
                st.success("Materiales integrados correctamente.")
                st.rerun()

    st.markdown("---")
    st.subheader("📚 Materiales Guardados")
    
    if not db["material"]:
        st.info("No hay registros en material.")
    else:
        for idx, item in enumerate(db["material"]):
            titulo_item = item.get("titulo") or item.get("titulo_registro") or "Sin Título"
            nombre_arc = item.get("nombre_original") or item.get("filename") or "Archivo"
            tipo_arc = item.get("tipo", "").lower()
            autor_arc = item.get("autor", "Desconocido")
            fecha_arc = item.get("fecha", "")
            resumen_arc = item.get("resumen", "Sin resumen disponible.")
            ruta_arc = item.get("ruta", "")
            
            with st.container():
                col_head, col_del_btn = st.columns([5, 1])
                with col_head:
                    st.markdown(f"### 📁 {titulo_item} — <span class='highlight-tag'>{nombre_arc}</span>", unsafe_allow_html=True)
                    st.caption(f"Autor: **{autor_arc}** | Fecha: {fecha_arc} | Formato: **{tipo_arc.upper()}**")
                with col_del_btn:
                    if usr.get("permiso_eliminar") or es_admin:
                        if st.button("🗑️ Borrar", key=f"del_mat_{idx}", type="secondary"):
                            if ruta_arc and os.path.exists(ruta_arc):
                                os.remove(ruta_arc)
                            if os.path.exists(os.path.join(DIR_USUARIOS, f"{nombre_arc}.json")):
                                os.remove(os.path.join(DIR_USUARIOS, f"{nombre_arc}.json"))
                            db["material"].pop(idx)
                            guardar_estado(db)
                            st.success("Archivo eliminado.")
                            st.rerun()
                
                t_vis, t_res = st.tabs(["👁️ Multimedia / Descarga", "📝 Diagnóstico y Resumen"])
                with t_vis:
                    if ruta_arc and os.path.exists(ruta_arc):
                        if tipo_arc == "mp4":
                            st.video(ruta_arc)
                        elif tipo_arc in ["m4a", "mp3", "wav"]:
                            st.audio(ruta_arc)
                        elif tipo_arc in ["jpg", "jpeg", "png"]:
                            st.image(ruta_arc, use_container_width=True)
                        else:
                            with open(ruta_arc, "rb") as fl:
                                st.download_button("📥 Descargar Archivo", data=fl.read(), file_name=nombre_arc, key=f"dl_a_{idx}")
                    else:
                        st.error("Archivo físico no encontrado.")
                with t_res:
                    st.markdown(resumen_arc)
                st.markdown("---")

# ==========================================
# 3. INFORMES COMPARTIDOS
# ==========================================
with pestanas_principales[2]:
    st.markdown("""<div class='modern-card'><h3 style='margin:0;'>📄 Informes Compartidos</h3><p style='margin:0; color:#4c1d95; font-weight:600;'>Generación y repositorio colaborativo con exportación en formato Carta.</p></div>""", unsafe_allow_html=True)
    
    with st.expander("🤖 Redactar Nuevo Informe con IA (Groq)", expanded=False):
        nom_i = st.text_input("Título del Informe:")
        enf_i = st.selectbox("Enfoque:", ["Resumen Ejecutivo", "Diagnóstico Técnico", "Evaluación Estratégica"])
        ins_i = st.text_area("Instrucciones complementarias:")
        
        if st.button("🚀 Generar Informe"):
            client = obtener_cliente_ia()
            if not nom_i.strip():
                st.error("Ingrese un título para el informe.")
            elif not client:
                st.error("API Key de Groq (GROQ_API_KEY) no configurada en Secrets.")
            else:
                with st.spinner("Redactando informe consolidado con Groq IA..."):
                    resumenes_mat = "\n".join([f"- {x.get('nombre_original', 'Archivo')}: {x.get('resumen', '')}" for x in db["material"]])
                    
                    p_sys = "Actúa como especialista analítico senior. Genera informes estructurados de excelencia en español latinoamericano."
                    p_user = f"""Genera un informe estructurado profesional con enfoque '{enf_i}'.
Título: {nom_i}
Autor solicitante: {usr['nombre']}
Instrucciones: {ins_i}
Materiales analizados: {resumenes_mat if resumenes_mat else 'Sin materiales adjuntos.'}

ESTRUCTURA REQUERIDA:
1. Diagnóstico y Visión Global
2. Hallazgos Analíticos Relevantes (Si incluyes tablas comparativas o síntesis, hazlo OBLIGATORIAMENTE usando sintaxis estándar de tablas Markdown con '|' y encabezados separados por '|---|---|', NUNCA uses arte ASCII ni bloques de código con ``` para tablas).
3. Conclusiones y Plan de Acción

REGLA OBLIGATORIA DE FORMATO:
- Las tablas deben ser tablas Markdown renderizables reales.
- Todo el texto analítico debe fluir con negritas y viñetas claras."""
                    try:
                        resp_txt = ejecutar_chat_groq(client, p_sys, p_user)
                        
                        nuevo_inf = {
                            "titulo": nom_i,
                            "autor": usr["nombre"],
                            "fecha": datetime.now().strftime("%Y-%m-%d %H:%M"),
                            "contenido": resp_txt
                        }
                        
                        db["informes"].append(nuevo_inf)
                        
                        nom_inf_f = f"inf_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
                        with open(os.path.join(DIR_INFORMES, nom_inf_f), "w", encoding="utf-8") as fi_b:
                            json.dump(nuevo_inf, fi_b, ensure_ascii=False, indent=2)
                            
                        guardar_estado(db)
                        st.success("Informe publicado y respaldado con éxito.")
                        st.rerun()
                    except Exception as err:
                        st.error(f"Error al generar informe: {err}")

    st.markdown("---")
    st.subheader("📚 Repositorio de Informes")
    
    if not db["informes"]:
        st.info("No hay informes registrados.")
    else:
        for idx_inf, inf in enumerate(db["informes"]):
            titulo_inf = inf.get("titulo") or inf.get("titulo_informe") or "Sin Título"
            autor_inf = inf.get("autor", "Desconocido")
            fecha_inf = inf.get("fecha", "")
            contenido_inf = inf.get("contenido", "")
            
            with st.container():
                c_inf_t, c_inf_del = st.columns([5, 1])
                with c_inf_t:
                    st.markdown(f"### 📄 [{autor_inf}] - {titulo_inf}")
                    st.caption(f"Generado el: {fecha_inf}")
                with c_inf_del:
                    if usr.get("permiso_eliminar") or es_admin:
                        if st.button("🗑️ Borrar", key=f"del_inf_{idx_inf}", type="secondary"):
                            db["informes"].pop(idx_inf)
                            guardar_estado(db)
                            st.success("Informe eliminado.")
                            st.rerun()
                
                st.markdown(
